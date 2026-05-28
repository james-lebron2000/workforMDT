"""报告导出 - Word / PDF / PPT / 微信卡片

依赖:
- python-docx (Word)
- weasyprint (HTML → PDF)
- python-pptx (PPT)

所有产物都带 "AI 辅助生成,需主治医师复核" 水印。
"""
from __future__ import annotations

import io
from datetime import date as _date
from typing import Any, Dict, List, Tuple

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from config import settings
from models.opinion import DepartmentOpinion
from models.patient import Patient
from models.recommendation import FinalRecommendation
from models.session import MdtSession
from models.summary import CaseSummary
from models.tnm import TnmStaging
from utils.logger import get_logger

logger = get_logger("export")

WATERMARK = "AI 辅助生成,需主治医师复核"


async def _load_bundle(db: AsyncSession, session_id: str) -> Dict[str, Any]:
    sess = await db.get(MdtSession, session_id)
    if sess is None:
        raise ValueError("session not found")
    patient = await db.get(Patient, sess.patient_id)
    summary = (
        await db.execute(
            select(CaseSummary)
            .where(CaseSummary.session_id == session_id)
            .order_by(CaseSummary.version.desc())
        )
    ).scalar_one_or_none()
    tnm = (
        await db.execute(
            select(TnmStaging)
            .where(TnmStaging.session_id == session_id)
            .order_by(TnmStaging.version.desc())
        )
    ).scalar_one_or_none()
    opinions = list(
        (
            await db.execute(
                select(DepartmentOpinion).where(DepartmentOpinion.session_id == session_id)
            )
        ).scalars()
    )
    final = (
        await db.execute(
            select(FinalRecommendation)
            .where(FinalRecommendation.session_id == session_id)
            .order_by(FinalRecommendation.version.desc())
        )
    ).scalar_one_or_none()
    return {
        "session": sess,
        "patient": patient,
        "summary": summary,
        "tnm": tnm,
        "opinions": opinions,
        "final": final,
    }


# ---------- Word ----------


def _build_docx(bundle: Dict[str, Any]) -> bytes:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    doc = Document()

    # 每页页眉 + 页脚加水印 —— python-docx 的 section.header/footer 会复制到每页,
    # 比只在首页放水印更难撕掉,符合"水印不可移除"的产品红线。
    section = doc.sections[0]
    hp = section.header.paragraphs[0]
    hp.text = ""
    hr = hp.add_run(f"[{WATERMARK}]")
    hr.italic = True
    hr.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
    hr.font.size = Pt(9)
    hp.alignment = WD_ALIGN_PARAGRAPH.CENTER

    fp = section.footer.paragraphs[0]
    fp.text = ""
    fr = fp.add_run(f"— {WATERMARK} · {settings.app_name} —")
    fr.italic = True
    fr.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    fr.font.size = Pt(8)
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER

    title = doc.add_heading(f"{settings.app_name} - 多学科病例整理报告", 0)
    for run in title.runs:
        run.bold = True

    p = doc.add_paragraph()
    r = p.add_run(f"[{WATERMARK}]")
    r.italic = True
    r.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)

    # QC 警告横幅(若 qc_status 是 warning/failed)
    final = bundle.get("final")
    if final and getattr(final, "qc_status", None) in ("warning", "failed"):
        wp = doc.add_paragraph()
        wr = wp.add_run(
            f"⚠ QC 状态:{final.qc_status}。本报告存在待处理问题,请在确认页修复后再分发。"
        )
        wr.bold = True
        wr.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)

    sess = bundle["session"]
    patient = bundle["patient"]
    doc.add_paragraph(
        f"会议时间:{sess.mdt_date or _date.today()}    "
        f"患者代号:{patient.code if patient else '-'}"
    )

    # 一、患者基本信息
    doc.add_heading("一、患者基本信息", level=1)
    if patient:
        doc.add_paragraph(
            f"性别:{patient.sex or '-'};年龄段:{patient.age_range or '-'};"
            f"主要诊断:{patient.primary_diagnosis or '-'};"
            f"原发部位:{patient.primary_site or '-'}"
        )

    # 二、本次就诊需求
    doc.add_heading("二、本次就诊需求与预期收获", level=1)
    summary = bundle["summary"]
    if summary:
        doc.add_paragraph(summary.chief_need or "-")
    else:
        doc.add_paragraph("-")

    # 三、病历摘要
    doc.add_heading("三、病历摘要", level=1)
    if summary:
        doc.add_paragraph(summary.history_summary or "-")

    # 四、既往治疗时间轴
    doc.add_heading("四、既往治疗时间轴", level=1)
    if summary and summary.treatment_timeline:
        for evt in summary.treatment_timeline:
            doc.add_paragraph(
                f"{evt.get('date','?')}:{evt.get('event','-')}",
                style="List Bullet",
            )

    # 五、当前临床判断
    doc.add_heading("五、当前临床判断", level=1)
    final = bundle["final"]
    if final and final.clinical_judgment:
        doc.add_paragraph(final.clinical_judgment)

    # 六、TNM 分期
    doc.add_heading("六、TNM 分期及依据", level=1)
    tnm = bundle["tnm"]
    if tnm:
        doc.add_paragraph(
            f"分期类型:{tnm.tnm_type};T={tnm.t_stage};N={tnm.n_stage};"
            f"M={tnm.m_stage};总分期:{tnm.overall_stage};置信度:{tnm.confidence}"
        )
        doc.add_paragraph(f"依据:{tnm.basis or '-'}")
        if tnm.uncertainty:
            doc.add_paragraph(f"不确定项:{tnm.uncertainty}")

    # 七、多学科医生意见
    doc.add_heading("七、多学科医生意见", level=1)
    for op in bundle["opinions"]:
        doc.add_heading(op.department, level=2)
        if op.is_missing:
            p = doc.add_paragraph("本次讨论未明确记录,建议补充")
            for r in p.runs:
                r.font.color.rgb = RGBColor(0xCC, 0x66, 0x00)
            continue
        if op.opinion:
            doc.add_paragraph(f"观点:{op.opinion}")
        if op.rationale:
            doc.add_paragraph(f"理由:{op.rationale}")
        if op.recommendation:
            doc.add_paragraph(f"建议:{op.recommendation}")
        if op.evidence_snippet:
            doc.add_paragraph(f"证据:{op.evidence_snippet}").italic = True

    # 八、检查建议
    doc.add_heading("八、检查建议", level=1)
    if final and final.exam_recommendations:
        for ex in final.exam_recommendations:
            doc.add_paragraph(
                f"[{ex.get('priority','建议')}] {ex.get('name','-')} - {ex.get('reason','')}",
                style="List Bullet",
            )

    # 九、治疗建议(需医生最终确认)
    doc.add_heading("九、治疗建议(需医生最终确认)", level=1)
    if final and final.treatment_recommendations:
        for t in final.treatment_recommendations:
            doc.add_paragraph(
                f"[{t.get('kind','-')}] {t.get('regimen','-')} - {t.get('rationale','')} "
                f"(证据等级 {t.get('evidence_level','未分级')})",
                style="List Bullet",
            )

    # 十、推荐就诊
    doc.add_heading("十、推荐就诊医生/门诊", level=1)
    if final and final.referral:
        for r in final.referral:
            line = f"{r.get('dept','-')}({r.get('doctor_hint','-')}) - {r.get('reason','')}"
            if r.get("bring_with"):
                line += f"\n  需携带:{', '.join(r['bring_with'])}"
            doc.add_paragraph(line, style="List Bullet")

    # 十一、患者反馈话术
    doc.add_heading("十一、给患者及家属的反馈话术", level=1)
    if final and final.patient_script:
        doc.add_paragraph(final.patient_script)

    # 末尾水印
    doc.add_paragraph()
    p = doc.add_paragraph()
    r = p.add_run(f"— {WATERMARK} —")
    r.italic = True
    r.font.size = Pt(9)

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


# ---------- Pre-MDT 摘要 PDF(会前阅读版,不含 TNM/科室意见/治疗建议) ----------


def _build_brief_html(bundle: Dict[str, Any]) -> str:
    """MDT 会前摘要 PDF — 只输出已与患者核对过的:
    诉求 + 病史 + 治疗时间轴 + MDT 待解答问题。

    刻意不输出 TNM / 多学科意见 / 治疗建议,因为这些是 MDT 中/后才产生,
    会前分享会误导与会医生。
    """
    sess = bundle["session"]
    patient = bundle["patient"]
    summary = bundle["summary"]

    def esc(s: Any) -> str:
        if s is None:
            return "-"
        return str(s).replace("<", "&lt;").replace(">", "&gt;")

    timeline_html = ""
    if summary and summary.treatment_timeline:
        timeline_html = "<ul>" + "".join(
            f"<li><b>{esc(e.get('date'))}</b>:{esc(e.get('event'))}</li>"
            for e in summary.treatment_timeline
        ) + "</ul>"

    questions_html = ""
    if summary and summary.mdt_questions:
        questions_html = "<ol>" + "".join(
            f"<li>{esc(q)}</li>" for q in summary.mdt_questions
        ) + "</ol>"

    # 路由层已硬门禁 (status >= summary_confirmed),走到这里就一定已确认;
    # 时间戳取 session.updated_at(confirm_summary 写状态时自动更新)。
    confirmed_at = (
        f"病史已与患者核对确认 · {sess.updated_at:%Y-%m-%d %H:%M}"
        if getattr(sess, "updated_at", None)
        else "病史已与患者核对确认"
    )

    return f"""<!doctype html>
<html><head><meta charset="utf-8"/>
<style>
  @page {{
    size: A4; margin: 22mm 16mm 22mm 16mm;
    @top-center {{ content: "[{WATERMARK} · 会前阅读版]"; color: #c00; font-style: italic; font-size: 9pt; }}
    @bottom-center {{ content: "— {WATERMARK} · MDT 会前阅读版 · 第 " counter(page) " / " counter(pages) " 页 —"; color: #999; font-size: 9pt; }}
  }}
  body {{ font-family: "Noto Sans CJK SC", "PingFang SC", sans-serif; color: #222; line-height: 1.55; }}
  h1 {{ color: #003366; border-bottom: 2px solid #003366; padding-bottom: 4px; }}
  h2 {{ color: #003366; margin-top: 22px; }}
  .watermark {{ color: #c00; font-style: italic; font-weight: bold; }}
  .confirm {{ color: #0a7; font-size: 11pt; }}
  .meta {{ color: #555; }}
  .notice {{ background: #fff5e6; border-left: 3px solid #f80; padding: 8px 12px; margin: 12px 0; font-size: 11pt; }}
  ul, ol {{ padding-left: 22px; }}
</style></head><body>
<h1>{esc(settings.app_name)} - MDT 会前病例摘要</h1>
<p class="watermark">[{WATERMARK}]</p>
<p class="confirm">{esc(confirmed_at)}</p>
<p class="meta">会议时间:{esc(sess.mdt_date or _date.today())} · 患者代号:{esc(patient.code) if patient else '-'}</p>

<div class="notice">
本摘要供与会医生<b>会前阅读</b>,基于已上传的资料(化验/影像/病理报告)由 AI 整理,
经协调医生与患者本人面对面核对。<br/>
本摘要<b>不包含</b>分期判断、各科意见或治疗建议 — 这些在 MDT 讨论中生成。
</div>

<h2>一、患者基本信息</h2>
<p>性别:{esc(patient.sex) if patient else '-'};年龄段:{esc(patient.age_range) if patient else '-'};
   主要诊断:{esc(patient.primary_diagnosis) if patient else '-'};
   原发部位:{esc(patient.primary_site) if patient else '-'}</p>

<h2>二、本次就诊需求与预期</h2>
<p>{esc(summary.chief_need if summary else '-')}</p>

<h2>三、病史摘要</h2>
<p>{esc(summary.history_summary if summary else '-')}</p>

<h2>四、既往治疗时间轴</h2>
{timeline_html or '<p>-</p>'}

<h2>五、当前问题</h2>
<p>{esc(summary.current_problem if summary else '-')}</p>

<h2>六、本次 MDT 待解答的问题</h2>
{questions_html or '<p>(暂未列出,请协调医生在会前补充)</p>'}

</body></html>"""


def _build_brief_pdf(bundle: Dict[str, Any]) -> bytes:
    from weasyprint import HTML

    html_str = _build_brief_html(bundle)
    return HTML(string=html_str).write_pdf()


# ---------- PDF (HTML 模板 + weasyprint) ----------


def _build_html(bundle: Dict[str, Any]) -> str:
    sess = bundle["session"]
    patient = bundle["patient"]
    summary = bundle["summary"]
    tnm = bundle["tnm"]
    final = bundle["final"]
    opinions = bundle["opinions"]

    def esc(s: Any) -> str:
        if s is None:
            return "-"
        return str(s).replace("<", "&lt;").replace(">", "&gt;")

    timeline_html = ""
    if summary and summary.treatment_timeline:
        timeline_html = "<ul>" + "".join(
            f"<li><b>{esc(e.get('date'))}</b>:{esc(e.get('event'))}</li>"
            for e in summary.treatment_timeline
        ) + "</ul>"

    opinions_html = ""
    for op in opinions:
        if op.is_missing:
            opinions_html += (
                f"<h3>{esc(op.department)}</h3>"
                f"<p class='missing'>本次讨论未明确记录,建议补充</p>"
            )
            continue
        opinions_html += (
            f"<h3>{esc(op.department)}</h3>"
            f"<p><b>观点</b>:{esc(op.opinion)}</p>"
            f"<p><b>理由</b>:{esc(op.rationale)}</p>"
            f"<p><b>建议</b>:{esc(op.recommendation)}</p>"
            + (f"<p class='ev'>证据:{esc(op.evidence_snippet)}</p>" if op.evidence_snippet else "")
        )

    exams_html = ""
    if final and final.exam_recommendations:
        exams_html = "<ul>" + "".join(
            f"<li>[{esc(e.get('priority','建议'))}] {esc(e.get('name'))} - {esc(e.get('reason'))}</li>"
            for e in final.exam_recommendations
        ) + "</ul>"

    treatments_html = ""
    if final and final.treatment_recommendations:
        treatments_html = "<ul>" + "".join(
            f"<li>[{esc(t.get('kind'))}] {esc(t.get('regimen'))} - {esc(t.get('rationale'))} "
            f"(证据等级 {esc(t.get('evidence_level','未分级'))})</li>"
            for t in final.treatment_recommendations
        ) + "</ul>"

    referral_html = ""
    if final and final.referral:
        referral_html = "<ul>" + "".join(
            f"<li>{esc(r.get('dept'))}({esc(r.get('doctor_hint'))}) - {esc(r.get('reason'))}"
            + (f"<br/>需携带:{esc(', '.join(r.get('bring_with',[]) ))}" if r.get('bring_with') else "")
            + "</li>"
            for r in final.referral
        ) + "</ul>"

    qc_status = getattr(final, "qc_status", None) if final else None
    qc_banner = ""
    if qc_status in ("warning", "failed"):
        qc_banner = (
            f'<p class="qc-bad"><b>⚠ QC 状态:{qc_status}</b>。'
            f'本报告存在待处理问题,请在确认页修复后再分发。</p>'
        )

    return f"""<!doctype html>
<html><head><meta charset="utf-8"/>
<style>
  @page {{
    size: A4;
    margin: 22mm 16mm 22mm 16mm;
    @top-center {{ content: "[{WATERMARK}]"; color: #c00; font-style: italic; font-size: 9pt; }}
    @bottom-center {{ content: "— {WATERMARK} · {esc(settings.app_name)} · 第 " counter(page) " 页 共 " counter(pages) " 页 —"; color: #999; font-size: 9pt; }}
  }}
  body {{ font-family: "Noto Sans CJK SC", "PingFang SC", sans-serif; color: #222; line-height: 1.5; }}
  h1 {{ color: #003366; border-bottom: 2px solid #003366; padding-bottom: 4px; }}
  h2 {{ color: #003366; margin-top: 22px; }}
  h3 {{ color: #555; margin-top: 12px; }}
  .watermark {{ color: #c00; font-style: italic; font-weight: bold; }}
  .qc-bad {{ color: #c00; background: #fff0f0; border-left: 3px solid #c00; padding: 6px 10px; margin: 8px 0; }}
  .missing {{ color: #c66; font-style: italic; }}
  .ev {{ color: #777; font-size: 11pt; font-style: italic; }}
  .meta {{ color: #555; }}
  ul {{ padding-left: 22px; }}
</style></head><body>
<h1>{esc(settings.app_name)} - 多学科病例整理报告</h1>
<p class="watermark">[{WATERMARK}]</p>
{qc_banner}
<p class="meta">会议时间:{esc(sess.mdt_date or _date.today())}    患者代号:{esc(patient.code) if patient else '-'}</p>

<h2>一、患者基本信息</h2>
<p>性别:{esc(patient.sex) if patient else '-'};年龄段:{esc(patient.age_range) if patient else '-'};
   主要诊断:{esc(patient.primary_diagnosis) if patient else '-'};
   原发部位:{esc(patient.primary_site) if patient else '-'}</p>

<h2>二、本次就诊需求与预期收获</h2>
<p>{esc(summary.chief_need if summary else '-')}</p>

<h2>三、病历摘要</h2>
<p>{esc(summary.history_summary if summary else '-')}</p>

<h2>四、既往治疗时间轴</h2>
{timeline_html or '<p>-</p>'}

<h2>五、当前临床判断</h2>
<p>{esc(final.clinical_judgment if final else '-')}</p>

<h2>六、TNM 分期及依据</h2>
<p>分期类型:{esc(tnm.tnm_type) if tnm else '-'};T={esc(tnm.t_stage) if tnm else '-'};
   N={esc(tnm.n_stage) if tnm else '-'};M={esc(tnm.m_stage) if tnm else '-'};
   总分期:{esc(tnm.overall_stage) if tnm else '-'};置信度:{esc(tnm.confidence) if tnm else '-'}</p>
<p><b>依据</b>:{esc(tnm.basis if tnm else '-')}</p>
<p><b>不确定项</b>:{esc(tnm.uncertainty if tnm else '-')}</p>

<h2>七、多学科医生意见</h2>
{opinions_html or '<p>-</p>'}

<h2>八、检查建议</h2>
{exams_html or '<p>-</p>'}

<h2>九、治疗建议(需医生最终确认)</h2>
{treatments_html or '<p>-</p>'}

<h2>十、推荐就诊医生/门诊</h2>
{referral_html or '<p>-</p>'}

<h2>十一、给患者及家属的反馈话术</h2>
<p>{esc(final.patient_script if final else '-')}</p>

</body></html>"""


def _build_pdf(bundle: Dict[str, Any]) -> bytes:
    from weasyprint import HTML

    html_str = _build_html(bundle)
    return HTML(string=html_str).write_pdf()


# ---------- PPT ----------


def _build_pptx(bundle: Dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # 标题 + 空白

    def _add_watermark(slide) -> None:
        """在每张 slide 底部加一行水印 — 红色斜体不可移除(撕掉要逐页操作)"""
        wm = slide.shapes.add_textbox(Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.3))
        tf = wm.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"⚠ {WATERMARK}"
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)

    def add_slide(title: str, body_lines: List[str]):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        txbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.2))
        tf = txbox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
        _add_watermark(slide)

    sess = bundle["session"]
    patient = bundle["patient"]
    summary = bundle["summary"]
    tnm = bundle["tnm"]
    final = bundle["final"]

    # 封面
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = f"{settings.app_name} - MDT 报告"
    cover.placeholders[1].text = (
        f"患者:{patient.code if patient else '-'}\n"
        f"会议时间:{sess.mdt_date or _date.today()}\n"
        f"[{WATERMARK}]"
    )
    _add_watermark(cover)

    if summary:
        add_slide(
            "患者诉求 + 病历摘要",
            [
                f"诉求:{summary.chief_need or '-'}",
                f"摘要:{(summary.history_summary or '-')[:300]}",
            ],
        )

    if tnm:
        add_slide(
            "TNM 分期",
            [
                f"类型:{tnm.tnm_type}",
                f"T/N/M:{tnm.t_stage} / {tnm.n_stage} / {tnm.m_stage}",
                f"总分期:{tnm.overall_stage}",
                f"依据:{(tnm.basis or '-')[:300]}",
                f"不确定项:{tnm.uncertainty or '-'}",
            ],
        )

    if bundle["opinions"]:
        add_slide(
            "多学科意见",
            [
                (
                    f"[{op.department}] " +
                    ("(未明确记录)" if op.is_missing else (op.opinion or "-"))
                )
                for op in bundle["opinions"][:8]
            ],
        )

    if final:
        add_slide(
            "治疗建议(需医生最终确认)",
            [
                f"[{t.get('kind')}] {t.get('regimen','-')}"
                for t in (final.treatment_recommendations or [])[:6]
            ],
        )
        add_slide("患者反馈话术", [(final.patient_script or "-")[:600]])

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


# ---------- 微信卡片(纯文本) ----------


def _build_wechat_card(bundle: Dict[str, Any]) -> bytes:
    """微信卡片(纯文本) — 水印放首行+末行+每段标题前,微信复制粘贴会保留。

    临床场景:卡片会被转发到家属群、其他医生群,撕掉水印很容易,
    所以多处冗余,确保任何一段被复制出去都带水印。
    """
    summary = bundle["summary"]
    tnm = bundle["tnm"]
    final = bundle["final"]
    patient = bundle["patient"]

    WM = f"⚠️ {WATERMARK}"

    lines: List[str] = []
    lines.append(f"【{settings.app_name} · MDT 报告卡片】")
    lines.append(WM)
    lines.append("━" * 18)
    lines.append(f"患者代号:{patient.code if patient else '-'}")
    if summary:
        lines.append(f"诉求:{summary.chief_need or '-'}")
    if tnm:
        lines.append(
            f"分期:{tnm.tnm_type} {tnm.t_stage}{tnm.n_stage}{tnm.m_stage} ({tnm.overall_stage})"
        )
        if tnm.basis:
            lines.append(f"  依据:{tnm.basis[:100]}")
    if final and final.clinical_judgment:
        lines.append("")
        lines.append(f"【临床判断】 {WM}")
        lines.append(final.clinical_judgment[:200])
    if final and final.treatment_recommendations:
        lines.append("")
        lines.append(f"【治疗建议(需医生最终确认)】 {WM}")
        for t in final.treatment_recommendations[:3]:
            lines.append(f"  · [{t.get('kind')}] {t.get('regimen')}")
    if final and final.qc_status in ("warning", "failed"):
        lines.append("")
        lines.append(f"⚠️ QC 状态:{final.qc_status} — 本卡片存在待处理问题,请回到确认页修复后再分发")
    lines.append("")
    lines.append("━" * 18)
    lines.append(WM)
    lines.append("(完整报告请向协调医生索取 Word/PDF)")
    return "\n".join(lines).encode("utf-8")


# ---------- 入口 ----------


async def build_report(
    db: AsyncSession, session_id: str, format_: str
) -> Tuple[bytes, str, str]:
    bundle = await _load_bundle(db, session_id)
    sess = bundle["session"]
    fname_base = f"MDT-{bundle['patient'].code if bundle['patient'] else 'unknown'}-{sess.mdt_date or _date.today()}"
    if format_ == "docx":
        return _build_docx(bundle), (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ), f"{fname_base}.docx"
    if format_ == "pdf":
        return _build_pdf(bundle), "application/pdf", f"{fname_base}.pdf"
    if format_ == "pptx":
        return _build_pptx(bundle), (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ), f"{fname_base}.pptx"
    if format_ == "wechat_card":
        return _build_wechat_card(bundle), "text/plain", f"{fname_base}.txt"
    if format_ == "brief_pdf":
        return _build_brief_pdf(bundle), "application/pdf", f"{fname_base}-brief.pdf"
    raise ValueError(f"unsupported format: {format_}")
